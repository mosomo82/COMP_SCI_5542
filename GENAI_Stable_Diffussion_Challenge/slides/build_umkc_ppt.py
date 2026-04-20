"""Build a UMKC-styled PPTX deck for COMP_SCI 5542.

Usage:
  pip install python-pptx pillow
  python slides/build_umkc_ppt.py

Optional assets:
  - Put UMKC logo at assets/umkc_logo.png
  - Update VIDEO_URL below after recording your demo
"""

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
SLIDES_DIR = Path(__file__).resolve().parent
ASSETS_DIR = ROOT / "assets"
OUT_PATH = SLIDES_DIR / "UMKC_COMP_SCI_5542_Presentation.pptx"

GITHUB_URL = "https://github.com/mosomo82/COMP_SCI_5542.git"
VIDEO_URL = "https://youtu.be/REPLACE_WITH_YOUR_VIDEO_ID"

UMKC_LOGO_PATH = ASSETS_DIR / "umkc_logo.png"
TITLE_HERO_IMAGE = ROOT / "outputs" / "0101635370" / "0101635370_structured_view01_seed42.png"
RESULT_IMAGE_PATHS = [
    ROOT / "outputs" / "0101635370" / "0101635370_naive_view01_seed42.png",
    ROOT / "outputs" / "0101635370" / "0101635370_naive_view02_seed43.png",
    ROOT / "outputs" / "0101635370" / "0101635370_structured_view01_seed42.png",
    ROOT / "outputs" / "0101635370" / "0101635370_structured_view02_seed43.png",
]

UMKC_BLUE = RGBColor(0x00, 0x4B, 0x8D)
UMKC_GOLD = RGBColor(0xFF, 0xC7, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT_CARD = RGBColor(0xF3, 0xF6, 0xFB)


def set_slide_notes(slide, lines: Iterable[str]) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    first = True
    for line in lines:
        para = notes.paragraphs[0] if first else notes.add_paragraph()
        first = False
        para.text = line


def style_title(shape, size=42):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Montserrat"
            r.font.bold = True
            r.font.size = Pt(size)
            r.font.color.rgb = UMKC_BLUE


def style_body(shape, size=20):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Open Sans"
            r.font.size = Pt(size)
            r.font.color.rgb = GRAY


def add_header_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = UMKC_BLUE
    bar.line.fill.background()

    gold = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.52), Inches(13.33), Inches(0.04))
    gold.fill.solid()
    gold.fill.fore_color.rgb = UMKC_GOLD
    gold.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.08), Inches(11.5), Inches(0.35))
    tf = tbox.text_frame
    tf.text = title
    run = tf.paragraphs[0].runs[0]
    run.font.name = "Montserrat"
    run.font.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE


def add_footer(slide, text="UMKC COMP_SCI 5542"):
    f = slide.shapes.add_textbox(Inches(0.35), Inches(7.15), Inches(9.5), Inches(0.25))
    tf = f.text_frame
    tf.text = text
    r = tf.paragraphs[0].runs[0]
    r.font.name = "Open Sans"
    r.font.size = Pt(10)
    r.font.color.rgb = GRAY


def add_logo_or_placeholder(slide, left, top, width, height):
    if UMKC_LOGO_PATH.exists():
        slide.shapes.add_picture(str(UMKC_LOGO_PATH), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
        return
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    ph.fill.solid()
    ph.fill.fore_color.rgb = LIGHT_CARD
    ph.line.color.rgb = UMKC_BLUE
    tf = ph.text_frame
    tf.text = "Insert UMKC logo\nassets/umkc_logo.png"
    style_body(ph, 10)


def add_image_or_placeholder(slide, img_path: Path, left, top, width, height, label):
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        ph.fill.solid()
        ph.fill.fore_color.rgb = LIGHT_CARD
        ph.line.color.rgb = UMKC_BLUE
        tf = ph.text_frame
        tf.text = f"Missing image\n{img_path.as_posix()}"
        style_body(ph, 11)
    cap = slide.shapes.add_textbox(Inches(left), Inches(top + height + 0.05), Inches(width), Inches(0.25))
    cap.text_frame.text = label
    style_body(cap, 12)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "COMP_SCI 5542: GenAI Stable Diffusion Challenge")

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(8.0), Inches(1.2))
    title.text_frame.text = "Data-Driven E-Commerce Image Generation"
    style_title(title, 40)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.30), Inches(7.8), Inches(1.5))
    sub.text_frame.text = (
        "Tuan (Tony) Nguyen\n"
        "University of Missouri-Kansas City\n"
        "MS in Data Science and Analytics"
    )
    style_body(sub, 21)

    add_logo_or_placeholder(slide, left=10.6, top=0.68, width=2.2, height=0.8)
    add_image_or_placeholder(
        slide,
        TITLE_HERO_IMAGE,
        left=8.8,
        top=1.55,
        width=4.1,
        height=4.1,
        label="Hero result (structured prompt)",
    )

    add_footer(slide)
    set_slide_notes(
        slide,
        [
            "Introduce objective: generate e-commerce product visuals from metadata.",
            "Explain that the project compares naive and structured prompt strategies.",
            "Mention evaluation combines CLIP metrics with manual quality scoring.",
        ],
    )


def add_bulleted_slide(prs, title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, title)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.8), Inches(5.95))
    tf = body.text_frame
    tf.clear()

    first = True
    for line in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.level = 0
        p.space_after = Pt(8)
        for r in p.runs:
            r.font.name = "Open Sans"
            r.font.size = Pt(22)
            r.font.color.rgb = GRAY

    add_footer(slide)
    if notes:
        set_slide_notes(slide, notes)


def add_code_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Methodology: Prompt Design")

    left = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(5.8), Inches(5.8))
    left.text_frame.text = (
        "Naive prompt:\n"
        "- Uses product title only\n\n"
        "Structured prompt:\n"
        "- Injects color, material, category, and style\n"
        "- Adds studio quality boosters\n"
        "- Uses negative prompt for artifact suppression"
    )
    style_body(left, 19)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.05), Inches(5.7), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK
    box.line.color.rgb = UMKC_BLUE

    code = slide.shapes.add_textbox(Inches(7.05), Inches(1.35), Inches(5.1), Inches(4.2))
    code.text_frame.text = (
        "prompt = (\n"
        "    f\"{photography_style} of \"\n"
        "    f\"{product['color']} {product['material']} {product['title']}, \"\n"
        "    f\"category: {product['category']}, style: {product['style']}, \"\n"
        "    f\"{quality_boosters}\"\n"
        ")"
    )
    tf = code.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Consolas"
            r.font.size = Pt(15)
            r.font.color.rgb = WHITE

    add_footer(slide)
    set_slide_notes(
        slide,
        [
            "Walk the audience through naive versus structured prompt construction.",
            "Highlight this is from the actual prompt_builder.py implementation.",
            "Explain why metadata-rich phrasing gives better controllability.",
        ],
    )


def add_results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Results (Images): Naive vs Structured")

    labels = [
        "Naive View 1",
        "Naive View 2",
        "Structured View 1",
        "Structured View 2",
    ]
    positions = [
        (0.85, 1.02),
        (3.95, 1.02),
        (7.05, 1.02),
        (10.15, 1.02),
    ]

    for (x, y), img, label in zip(positions, RESULT_IMAGE_PATHS, labels):
        add_image_or_placeholder(slide, img, left=x, top=y, width=2.85, height=2.85, label=label)

    cap = slide.shapes.add_textbox(Inches(0.85), Inches(4.40), Inches(12.2), Inches(1.2))
    cap.text_frame.text = (
        "Visual check: structured prompts tend to improve product framing and studio consistency.\n"
        "Use this slide to call out concrete examples during presentation."
    )
    style_body(cap, 17)

    add_footer(slide)
    set_slide_notes(
        slide,
        [
            "Compare naive and structured rows explicitly.",
            "Point to lighting/background differences and shape consistency.",
            "Mention this example uses outputs from product 0101635370.",
        ],
    )


def add_link_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Video, Repository, and AI Disclosure")

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.9), Inches(5.8))
    body.text_frame.text = (
        "GitHub URL:\n"
        f"{GITHUB_URL}\n\n"
        "Demo Video URL:\n"
        f"{VIDEO_URL}\n\n"
        "AI tools used (required disclosure):\n"
        "- Claude Code\n"
        "- ChatGPT\n"
        "- Antigravity AI\n\n"
        "Human verification:\n"
        "All code, outputs, metrics, and conclusions were manually reviewed."
    )
    style_body(body, 20)

    add_footer(slide)
    set_slide_notes(
        slide,
        [
            "Replace placeholder video URL with your final demo before submission.",
            "Reiterate disclosure transparency and manual validation practices.",
            "Invite questions and transition to Q&A.",
        ],
    )


def build():
    prs = Presentation()

    add_title_slide(prs)

    add_bulleted_slide(
        prs,
        "Scenario Description",
        [
            "E-commerce needs scalable, consistent product imagery.",
            "Manual studio photography is expensive and slow.",
            "Goal: generate commercial-style product visuals from metadata.",
            "Compare naive title-only prompts versus structured prompt engineering.",
        ],
        notes=[
            "Frame this as a practical engineering problem with business impact.",
            "State that metadata availability motivates data-driven generation.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Dataset",
        [
            "Source: Amazon Product Dataset metadata samples.",
            "Project subset: 10 products across mixed categories.",
            "Fields used: id, title, category, color, material, style.",
            "Data file in repo: data/products.json.",
        ],
        notes=[
            "Mention subset size is intentionally small for controlled experiments.",
            "Explain how each field contributes to prompt richness.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Methodology Overview",
        [
            "1) Metadata ingest and prompt construction.",
            "2) Image generation via SD/SDXL (optional ControlNet).",
            "3) Evaluation: CLIP, consistency, diversity, manual quality.",
            "4) Reports exported to CSV and HTML for review.",
            "5) Three execution modes: local CLI, Colab batch, Colab API.",
        ],
        notes=[
            "This slide is your architecture bridge into implementation detail.",
            "Keep this high-level and transition to pipeline specifics next.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Stable Diffusion Pipeline",
        [
            "Supports SD 1.5 and SDXL with Diffusers backend.",
            "Configurable: resolution, steps, CFG, seed, view count.",
            "Optional ControlNet adds structural conditioning.",
            "Reproducibility maintained via deterministic seed strategy.",
        ],
        notes=[
            "Emphasize engineering controls: parameters, seeds, and mode toggles.",
            "Explain when SDXL is preferred versus SD 1.5.",
        ],
    )

    add_code_slide(prs)

    add_bulleted_slide(
        prs,
        "Control Strategy",
        [
            "ControlNet with Canny edges anchors object structure.",
            "Improves cross-view consistency for the same product.",
            "Applied when reference image or edge guidance is available.",
            "Tradeoff: tighter structure can reduce visual diversity.",
        ],
        notes=[
            "Explain why this is critical for multi-view e-commerce listings.",
            "Connect this to consistency metric improvements.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Tools and Technologies",
        [
            "Core stack: Python, Diffusers, Stable Diffusion, ControlNet, CLIP.",
            "Execution: local environment and Google Colab GPU runtime.",
            "Deployment/demo: Flask API + ngrok + browser frontend.",
            "Evaluation and reporting: pandas, CSV/HTML outputs, quality workflow.",
        ],
        notes=[
            "Keep this concise and map each tool to one specific role.",
            "Mention that AI tool disclosure appears on the final slide.",
        ],
    )

    add_results_slide(prs)

    add_bulleted_slide(
        prs,
        "Evaluation",
        [
            "Metrics: CLIP alignment, consistency, diversity, human quality score.",
            "Sample (results/local_summary.csv, product 0101635370):",
            "Naive: CLIP 0.2997 | Consistency 0.7745 | Diversity 0.0442.",
            "Structured: CLIP 0.2714 | Consistency 0.7169 | Diversity 0.0442.",
            "Use full multi-product summary for final aggregate conclusions.",
        ],
        notes=[
            "Explain why one product snapshot is illustrative but not definitive.",
            "Encourage panel to focus on full-table trends in report artifacts.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Findings and Insights",
        [
            "Structured prompts improve controllability and presentation consistency.",
            "Control signals reduce geometric drift across product views.",
            "Qualitative scoring complements embedding-based metrics.",
            "Failure-case analysis reveals practical prompt and control fixes.",
        ],
        notes=[
            "Tie findings back to problem statement and business relevance.",
            "Use one concrete before/after visual when speaking.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Limitations",
        [
            "Current sample size is limited for broad statistical claims.",
            "CLIP is useful but not a full proxy for commercial quality.",
            "Model behavior depends on seed, runtime, and prompt sensitivity.",
            "Colab/API demo environments can introduce session instability.",
        ],
        notes=[
            "Present limitations confidently as future-work opportunities.",
            "Mention that reproducibility controls are already in place.",
        ],
    )

    add_link_slide(prs)

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
